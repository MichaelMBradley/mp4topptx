from cv2 import VideoCapture, CAP_PROP_FPS, imwrite
from os import remove, rmdir, mkdir
from os.path import exists, dirname, isfile
from glob import glob
from pptx import Presentation
from pptx.util import Inches
from sys import argv, exit
from status import statusbar, timer

# Ensuring valid argument
# Returns (path to video, path to frame directory)
def validate_input(arguments: list[str]) -> tuple[str, str]:
    if len(arguments) != 2:
        exit("Invalid argument (video path should be the only argument)")

    videofile = arguments[1]

    if videofile[-4:] != ".mp4":
        exit("Invalid format (format should be .mp4)")

    if not isfile(videofile):
        exit("File does not exist")

    frames = dirname(videofile) + "\\frames\\"

    return (videofile, frames)


# Function to set up path
def delete_frames(createdir: bool, frames: str) -> None:
    try:
        if exists(frames):
            # TODO: For whatever reason this isn't working
            for filename in glob(frames + "*"):
                print(filename, exists(filename))
                remove(filename)  # Deletes each file in \frames\
            rmdir(frames)
        if createdir:
            mkdir(frames)
    except OSError as error:
        exit(error)


# Returns (number of frames, frame rate)
def read_video(videofile: str, frames: str) -> tuple[int, int]:
    vid = VideoCapture(videofile)
    framerate = vid.get(CAP_PROP_FPS)

    # Reading video
    framenumber = 0
    while True:
        framesleft, frame = vid.read()
        if not framesleft:
            break
        else:
            # Writes each frame to file
            imwrite(f"{frames}{framenumber}.jpg", frame)
            framenumber += 1
    vid.release()

    return (framenumber, framerate)


def convert(arguments: list[str]) -> None:
    # Initializing
    (videofile, frames) = validate_input(arguments)
    delete_frames(True, frames)

    # Creating pptx
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    (totalframes, framerate) = read_video(videofile, frames)
    freq = framerate / 30
    framenumber = 0

    # Initialize time-keeping
    status = statusbar(int(totalframes / freq), "Creating pptx")
    timing = timer(["Slides", "Saving"])
    timing.start()

    # Create each slide
    while int(framenumber) < totalframes:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(f"{frames}{int(framenumber)}.jpg", 0, 0, width=Inches(16), height=Inches(9))
        framenumber += freq
        status.incrementandprint()
        timing.swapto(0)

    # Clean-up
    prs.save(f"{videofile[:-3]}pptx")
    timing.swapto(1)
    delete_frames(False, frames)
    timing.stop()
    timing.results()


if __name__ == "__main__":
    convert(argv)
